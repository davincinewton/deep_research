# This is copied from Magentic-one's great repo: https://github.com/microsoft/autogen/blob/v0.4.4/python/packages/autogen-magentic-one/src/autogen_magentic_one/markdown_browser/mdconvert.py
# Thanks to Microsoft researchers for open-sourcing this!
# type: ignore
import base64
import copy
import html
import json
import mimetypes
import os
import re
import shutil
import subprocess
import sys
import tempfile
import traceback
import zipfile
from typing import Any, Dict, List, Optional, Union
from urllib.parse import parse_qs, quote, unquote, urlparse, urlunparse

import mammoth
import markdownify
import pandas as pd
import fitz 
import pptx

# File-format detection
import puremagic
import pydub
import requests
import speech_recognition as sr
from bs4 import BeautifulSoup
from youtube_transcript_api import YouTubeTranscriptApi
from youtube_transcript_api.formatters import SRTFormatter
from web_vision_query_playwright_v4 import run_base64_analyser
from PIL import Image
from io import BytesIO
# from openai import OpenAI

from marker.converters.pdf import PdfConverter as pdfc
from marker.models import create_model_dict
from marker.output import text_from_rendered

class _CustomMarkdownify(markdownify.MarkdownConverter):
    """
    A custom version of markdownify's MarkdownConverter. Changes include:

    - Altering the default heading style to use '#', '##', etc.
    - Removing javascript hyperlinks.
    - Truncating images with large data:uri sources.
    - Ensuring URIs are properly escaped, and do not conflict with Markdown syntax
    """

    def __init__(self, **options: Any):
        options["heading_style"] = options.get("heading_style", markdownify.ATX)
        # Explicitly cast options to the expected type if necessary
        super().__init__(**options)

    def convert_hn(self, n: int, el: Any, text: str, convert_as_inline: bool) -> str:
        """Same as usual, but be sure to start with a new line"""
        if not convert_as_inline:
            if not re.search(r"^\n", text):
                return "\n" + super().convert_hn(n, el, text, convert_as_inline)  # type: ignore

        return super().convert_hn(n, el, text, convert_as_inline)  # type: ignore

    def convert_a(self, el: Any, text: str, convert_as_inline: bool):
        """Same as usual converter, but removes Javascript links and escapes URIs."""
        prefix, suffix, text = markdownify.chomp(text)  # type: ignore
        if not text:
            return ""
        href = el.get("href")
        title = el.get("title")

        # Escape URIs and skip non-http or file schemes
        if href:
            try:
                parsed_url = urlparse(href)  # type: ignore
                if parsed_url.scheme and parsed_url.scheme.lower() not in ["http", "https", "file"]:  # type: ignore
                    return "%s%s%s" % (prefix, text, suffix)
                href = urlunparse(parsed_url._replace(path=quote(unquote(parsed_url.path))))  # type: ignore
            except ValueError:  # It's not clear if this ever gets thrown
                return "%s%s%s" % (prefix, text, suffix)

        # For the replacement see #29: text nodes underscores are escaped
        if (
            self.options["autolinks"]
            and text.replace(r"\_", "_") == href
            and not title
            and not self.options["default_title"]
        ):
            # Shortcut syntax
            return "<%s>" % href
        if self.options["default_title"] and not title:
            title = href
        title_part = ' "%s"' % title.replace('"', r"\"") if title else ""
        return "%s[%s](%s%s)%s" % (prefix, text, href, title_part, suffix) if href else text

    def convert_img(self, el: Any, text: str, convert_as_inline: bool) -> str:
        """Same as usual converter, but removes data URIs"""

        alt = el.attrs.get("alt", None) or ""
        src = el.attrs.get("src", None) or ""
        title = el.attrs.get("title", None) or ""
        title_part = ' "%s"' % title.replace('"', r"\"") if title else ""
        if convert_as_inline and el.parent.name not in self.options["keep_inline_images_in"]:
            return alt

        # Remove dataURIs
        if src.startswith("data:"):
            src = src.split(",")[0] + "..."

        return "![%s](%s%s)" % (alt, src, title_part)

    def convert_soup(self, soup: Any) -> str:
        return super().convert_soup(soup)  # type: ignore


class DocumentConverterResult:
    """The result of converting a document to text."""

    def __init__(self, title: Union[str, None] = None, text_content: str = ""):
        self.title: Union[str, None] = title
        self.text_content: str = text_content


class DocumentConverter:
    """Abstract superclass of all DocumentConverters."""

    def convert(self, local_path: str, **kwargs: Any) -> Union[None, DocumentConverterResult]:
        raise NotImplementedError()


class PlainTextConverter(DocumentConverter):
    """Anything with content type text/plain"""

    def convert(self, local_path: str, **kwargs: Any) -> Union[None, DocumentConverterResult]:
        # Guess the content type from any file extension that might be around
        content_type, _ = mimetypes.guess_type("__placeholder" + kwargs.get("file_extension", ""))

        # Only accept text files
        if content_type is None:
            return None
        # elif "text/" not in content_type.lower():
        #     return None

        text_content = ""
        with open(local_path, "rt", encoding="utf-8") as fh:
            text_content = fh.read()
        return DocumentConverterResult(
            title=None,
            text_content=text_content,
        )


class HtmlConverter(DocumentConverter):
    """Anything with content type text/html"""

    def convert(self, local_path: str, **kwargs: Any) -> Union[None, DocumentConverterResult]:
        # Bail if not html
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".html", ".htm"]:
            return None

        result = None
        with open(local_path, "rt", encoding="utf-8") as fh:
            result = self._convert(fh.read())

        return result

    def _convert(self, html_content: str) -> Union[None, DocumentConverterResult]:
        """Helper function that converts and HTML string."""

        # Parse the string
        soup = BeautifulSoup(html_content, "html.parser")

        # Remove javascript and style blocks
        for script in soup(["script", "style"]):
            script.extract()

        # Print only the main content
        body_elm = soup.find("body")
        webpage_text = ""
        if body_elm:
            webpage_text = _CustomMarkdownify().convert_soup(body_elm)
        else:
            webpage_text = _CustomMarkdownify().convert_soup(soup)

        assert isinstance(webpage_text, str)

        return DocumentConverterResult(
            title=None if soup.title is None else soup.title.string, text_content=webpage_text
        )


class WikipediaConverter(DocumentConverter):
    """Handle Wikipedia pages separately, focusing only on the main document content."""

    def convert(self, local_path: str, **kwargs: Any) -> Union[None, DocumentConverterResult]:
        # Bail if not Wikipedia
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".html", ".htm"]:
            return None
        url = kwargs.get("url", "")
        if not re.search(r"^https?:\/\/[a-zA-Z]{2,3}\.wikipedia.org\/", url):
            return None

        # Parse the file
        soup = None
        with open(local_path, "rt", encoding="utf-8") as fh:
            soup = BeautifulSoup(fh.read(), "html.parser")

        # Remove javascript and style blocks
        for script in soup(["script", "style"]):
            script.extract()

        # Print only the main content
        body_elm = soup.find("div", {"id": "mw-content-text"})
        title_elm = soup.find("span", {"class": "mw-page-title-main"})

        webpage_text = ""
        main_title = None if soup.title is None else soup.title.string

        if body_elm:
            # What's the title
            if title_elm and len(title_elm) > 0:
                main_title = title_elm.string  # type: ignore
                assert isinstance(main_title, str)

            # Convert the page
            webpage_text = f"# {main_title}\n\n" + _CustomMarkdownify().convert_soup(body_elm)
        else:
            webpage_text = _CustomMarkdownify().convert_soup(soup)

        return DocumentConverterResult(
            title=main_title,
            text_content=webpage_text,
        )


class YouTubeConverter(DocumentConverter):
    """Handle YouTube specially, focusing on the video title, description, and transcript."""

    def convert(self, local_path: str, **kwargs: Any) -> Union[None, DocumentConverterResult]:
        # Bail if not YouTube
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".html", ".htm"]:
            return None
        url = kwargs.get("url", "")
        if not url.startswith("https://www.youtube.com/watch?"):
            return None

        # Parse the file
        soup = None
        with open(local_path, "rt", encoding="utf-8") as fh:
            soup = BeautifulSoup(fh.read(), "html.parser")

        # Read the meta tags
        assert soup.title is not None and soup.title.string is not None
        metadata: Dict[str, str] = {"title": soup.title.string}
        for meta in soup(["meta"]):
            for a in meta.attrs:
                if a in ["itemprop", "property", "name"]:
                    metadata[meta[a]] = meta.get("content", "")
                    break

        # We can also try to read the full description. This is more prone to breaking, since it reaches into the page implementation
        try:
            for script in soup(["script"]):
                content = script.text
                if "ytInitialData" in content:
                    lines = re.split(r"\r?\n", content)
                    obj_start = lines[0].find("{")
                    obj_end = lines[0].rfind("}")
                    if obj_start >= 0 and obj_end >= 0:
                        data = json.loads(lines[0][obj_start : obj_end + 1])
                        attrdesc = self._findKey(data, "attributedDescriptionBodyText")  # type: ignore
                        if attrdesc:
                            metadata["description"] = str(attrdesc["content"])
                    break
        except Exception:
            pass

        # Start preparing the page
        webpage_text = "# YouTube\n"

        title = self._get(metadata, ["title", "og:title", "name"])  # type: ignore
        assert isinstance(title, str)

        if title:
            webpage_text += f"\n## {title}\n"

        stats = ""
        views = self._get(metadata, ["interactionCount"])  # type: ignore
        if views:
            stats += f"- **Views:** {views}\n"

        keywords = self._get(metadata, ["keywords"])  # type: ignore
        if keywords:
            stats += f"- **Keywords:** {keywords}\n"

        runtime = self._get(metadata, ["duration"])  # type: ignore
        if runtime:
            stats += f"- **Runtime:** {runtime}\n"

        if len(stats) > 0:
            webpage_text += f"\n### Video Metadata\n{stats}\n"

        description = self._get(metadata, ["description", "og:description"])  # type: ignore
        if description:
            webpage_text += f"\n### Description\n{description}\n"

        transcript_text = ""
        parsed_url = urlparse(url)  # type: ignore
        params = parse_qs(parsed_url.query)  # type: ignore
        if "v" in params:
            assert isinstance(params["v"][0], str)
            video_id = str(params["v"][0])
            try:
                # Must be a single transcript.
                transcript = YouTubeTranscriptApi.get_transcript(video_id)  # type: ignore
                # transcript_text = " ".join([part["text"] for part in transcript])  # type: ignore
                # Alternative formatting:
                transcript_text = SRTFormatter().format_transcript(transcript)
            except Exception:
                pass
        if transcript_text:
            webpage_text += f"\n### Transcript\n{transcript_text}\n"

        title = title if title else soup.title.string
        assert isinstance(title, str)

        return DocumentConverterResult(
            title=title,
            text_content=webpage_text,
        )

    def _get(self, metadata: Dict[str, str], keys: List[str], default: Union[str, None] = None) -> Union[str, None]:
        for k in keys:
            if k in metadata:
                return metadata[k]
        return default

    def _findKey(self, json: Any, key: str) -> Union[str, None]:  # TODO: Fix json type
        if isinstance(json, list):
            for elm in json:
                ret = self._findKey(elm, key)
                if ret is not None:
                    return ret
        elif isinstance(json, dict):
            for k in json:
                if k == key:
                    return json[k]
                else:
                    ret = self._findKey(json[k], key)
                    if ret is not None:
                        return ret
        return None

def classify_pdf_page(page):
    """
    Classifies a single PDF page.
    Returns: "digital_native", "ocr_based", "image_only", "empty_page"
    """
    # Thresholds (can be tuned)
    MIN_CHARS_FOR_TEXT_PAGE = 20  # Minimum characters to consider a page as having text
    LARGE_IMAGE_AREA_THRESHOLD = 0.75  # If an image covers >75% of the page area
    MIN_TRACED_CHARS_FOR_OCR = 10 # Minimum characters detected by get_texttrace() to suspect OCR    
    page_area = page.rect.width * page.rect.height
    if page_area == 0: # Avoid division by zero for empty pages
        return "empty_page"

    # 1. Check for plain extractable text
    plain_text = page.get_text().strip()
    if len(plain_text) < MIN_CHARS_FOR_TEXT_PAGE:
        # If very little or no plain text, it might be image-only or very sparse
        # Check if there are any images.
        if not page.get_images(full=True):
             # No text and no images. Could be blank or vector-only.
             # For simplicity, if truly empty (no drawing commands either), it's "empty_page".
             # Otherwise, could be "digital_native" (e.g. vector drawings).
             # A more sophisticated check could look at drawing commands.
             # If page.get_drawings() is empty, it's truly blank.
             if not page.get_drawings():
                 return "empty_page"
             return "digital_native" # Or "vector_only"
        return "image_only"

    # 2. Check for "hidden" text (text rendering mode 3) or Type 3 font text
    # get_texttrace() returns a list of dictionaries.
    traced_text_objects = page.get_texttrace()
    num_traced_chars = 0
    if traced_text_objects: # Check if the list is not empty
        num_traced_chars = sum(len(trace_item.get('text', '')) # Use .get for safety
                               for trace_item in traced_text_objects
                               if trace_item.get('text')) # Ensure 'text' key exists and is not empty

    if num_traced_chars > MIN_TRACED_CHARS_FOR_OCR:
        # If a significant amount of text is "traced"
        # Check if there's a large image, common in OCR
        images = page.get_images(full=True)
        for img_info in images:
            try:
                img_rect = page.get_image_bbox(img_info)
                img_area = img_rect.width * img_rect.height if img_rect else 0
                if img_area / page_area > LARGE_IMAGE_AREA_THRESHOLD:
                    return "ocr_based" # Hidden text + large image = strong OCR
            except Exception: # Handle cases where bbox might not be available
                pass
        # If hidden text found but no large image, could still be OCR or complex digital
        return "ocr_based" # Lean towards OCR if significant traced text

    # 3. Check for large background images with normal text on top
    images = page.get_images(full=True)
    has_large_image = False
    for img_info in images:
        try:
            img_rect = page.get_image_bbox(img_info) # Get image bounding box
            img_area = img_rect.width * img_rect.height if img_rect else 0
            if img_area == 0 : continue
            if (img_area / page_area) > LARGE_IMAGE_AREA_THRESHOLD:
                has_large_image = True
                break
        except Exception: # Handle cases where bbox might not be available
            pass
    
    if has_large_image:
        # If there's text AND a large background image, it's suspicious.
        # This is a weaker signal than traced_text but still indicative of OCR.
        return "ocr_based"

    # 4. If text exists, no significant traced/hidden text, and no dominant large image
    return "digital_native"


def get_pdf_type(pdf_path):
    """
    Analyzes a PDF to determine if it's likely OCR-based, digital native, or image-only.
    Returns a dictionary with page_types and an overall_type.
    """
    try:
        doc = fitz.open(pdf_path)
    except Exception as e:
        return {"error": f"Could not open PDF: {e}"}

    page_classifications = []
    if not doc.page_count:
        return {"overall_type": "empty_pdf", "page_types": []}

    for i in range(min(doc.page_count, 3)): # Analyze first 5 pages for speed, or all
        page = doc.load_page(i)
        page_type = classify_pdf_page(page)
        page_classifications.append(page_type)

    doc.close()

    # Determine overall type
    overall_type = "unknown"
    if not page_classifications:
        overall_type = "empty_or_unreadable"
    elif all(pt == "ocr_based" for pt in page_classifications):
        overall_type = "ocr_based"
    elif all(pt == "digital_native" for pt in page_classifications):
        overall_type = "digital_native"
    elif all(pt == "image_only" for pt in page_classifications):
        overall_type = "image_only"
    elif "ocr_based" in page_classifications:
        overall_type = "mixed_contains_ocr" # If any page is OCR, flag it
    elif "digital_native" in page_classifications and "image_only" in page_classifications:
        overall_type = "mixed_digital_and_image"
    elif "digital_native" in page_classifications: # Only digital and empty/other
        overall_type = "digital_native"
    elif "image_only" in page_classifications: # Only image and empty/other
        overall_type = "image_only"
    else: # e.g. all empty_page
        overall_type = "mostly_empty_or_unclassified"


    return {
        "overall_type": overall_type,
        "analyzed_page_count": len(page_classifications),
        "page_types": page_classifications
    }

def ocr_pdf(pdf_path, cnt_pages=0):
    # Open the PDF
    doc = fitz.open(pdf_path)
    # images = []
    prompt = "Read all the text in the image, then out in markdown"
    # print(prompt)
    # client_x = OpenAI(
    #     api_key="lm-studio",
    #     base_url="http://127.0.0.1:1234/v1",
    # )
    # model_id_x="qwen2.5-vl-3b-instruct"
    # print(prompt)
    ocr_result = []
    # Convert each page to an image
    if cnt_pages<=0:
        cnt = len(doc)
    else:
        cnt = max(1,min(cnt_pages,len(doc)))
    for page_num in range(cnt):
        try:
            page = doc.load_page(page_num)
            pix = page.get_pixmap(matrix=fitz.Matrix(100/72, 100/72))  # 300 DPI
            print(pix.width, pix.height)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            # images.append(img)
            # Save image to bytes buffer
            # print('img done')
            buffer = BytesIO()
            img.save(buffer, format="JPEG")  # Save as PNG or JPEG
            img_bytes = buffer.getvalue()
            # print('buffer done')
            # Encode to Base64
            base64_string = base64.b64encode(img_bytes).decode('utf-8')
            content_type = "image/jpeg"
            data_uri = f"data:{content_type};base64,{base64_string}"
            print('image send to vlm...')
            # result = run_base64_analyser(prompt,data_uri,modelid=model_id_x,clientapi=client_x)
            result = run_base64_analyser(prompt,data_uri)
            ocr_result.append(result)
        except Exception as e:
            print(e)
    doc.close()    
    # print(ocr_result)
    return ocr_result

def digi_pdf(pdf_path):
    converter = pdfc(
        artifact_dict=create_model_dict(),
    )
    # rendered = converter(r"G:\360Downloads\021a5339-744f-42b7-bd9b-9368b3efda7a.pdf")
    rendered = converter(pdf_path)
    text, _, images = text_from_rendered(rendered)
    # print(text)
    # print(images)
    return text

class PdfConverter(DocumentConverter):
    """
    Converts PDFs to Markdown. Most style information is ignored, so the results are essentially plain-text.
    """
    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not a PDF
        extension = kwargs.get("file_extension", "")
        # print(local_path)
        # print(extension)
        if extension.lower() != ".pdf":

            return None
        # print('get in')
        result_digital = get_pdf_type(local_path)
        print(f"'{local_path}': {result_digital['overall_type']} (Pages: {result_digital['page_types']})")
        if result_digital['overall_type'] == "digital_native":
            return DocumentConverterResult(
                title=None,
                text_content=digi_pdf(local_path),
            )
        else:
            fulltext =  ocr_pdf(local_path)
            tmp = "\n\n".join(fulltext)
            # print(tmp)
            return DocumentConverterResult(
                title=None,
                text_content=tmp,
            )            


class DocxConverter(HtmlConverter):
    """
    Converts DOCX files to Markdown. Style information (e.g.m headings) and tables are preserved where possible.
    """

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not a DOCX
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".docx":
            return None

        result = None
        with open(local_path, "rb") as docx_file:
            result = mammoth.convert_to_html(docx_file)
            html_content = result.value
            result = self._convert(html_content)

        return result


class XlsxConverter(HtmlConverter):
    """
    Converts XLSX files to Markdown, with each sheet presented as a separate Markdown table.
    """

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not a XLSX
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".xlsx", ".xls"]:
            return None

        sheets = pd.read_excel(local_path, sheet_name=None)
        md_content = ""
        for s in sheets:
            md_content += f"## {s}\n"
            html_content = sheets[s].to_html(index=False)
            md_content += self._convert(html_content).text_content.strip() + "\n\n"

        return DocumentConverterResult(
            title=None,
            text_content=md_content.strip(),
        )


class PptxConverter(HtmlConverter):
    """
    Converts PPTX files to Markdown. Supports heading, tables and images with alt text.
    """

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not a PPTX
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".pptx":
            return None

        md_content = ""

        presentation = pptx.Presentation(local_path)
        slide_num = 0
        for slide in presentation.slides:
            slide_num += 1

            md_content += f"\n\n<!-- Slide number: {slide_num} -->\n"

            title = slide.shapes.title
            for shape in slide.shapes:
                # Pictures
                if self._is_picture(shape):
                    # https://github.com/scanny/python-pptx/pull/512#issuecomment-1713100069
                    alt_text = ""
                    try:
                        alt_text = shape._element._nvXxPr.cNvPr.attrib.get("descr", "")
                    except Exception:
                        pass

                    # A placeholder name
                    filename = re.sub(r"\W", "", shape.name) + ".jpg"
                    md_content += "\n![" + (alt_text if alt_text else shape.name) + "](" + filename + ")\n"

                # Tables
                if self._is_table(shape):
                    html_table = "<html><body><table>"
                    first_row = True
                    for row in shape.table.rows:
                        html_table += "<tr>"
                        for cell in row.cells:
                            if first_row:
                                html_table += "<th>" + html.escape(cell.text) + "</th>"
                            else:
                                html_table += "<td>" + html.escape(cell.text) + "</td>"
                        html_table += "</tr>"
                        first_row = False
                    html_table += "</table></body></html>"
                    md_content += "\n" + self._convert(html_table).text_content.strip() + "\n"

                # Text areas
                elif shape.has_text_frame:
                    if shape == title:
                        md_content += "# " + shape.text.lstrip() + "\n"
                    else:
                        md_content += shape.text + "\n"

            md_content = md_content.strip()

            if slide.has_notes_slide:
                md_content += "\n\n### Notes:\n"
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame is not None:
                    md_content += notes_frame.text
                md_content = md_content.strip()

        return DocumentConverterResult(
            title=None,
            text_content=md_content.strip(),
        )

    def _is_picture(self, shape):
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
            return True
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PLACEHOLDER:
            if hasattr(shape, "image"):
                return True
        return False

    def _is_table(self, shape):
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.TABLE:
            return True
        return False


class MediaConverter(DocumentConverter):
    """
    Abstract class for multi-modal media (e.g., images and audio)
    """

    def _get_metadata(self, local_path):
        exiftool = shutil.which("exiftool")
        if not exiftool:
            return None
        else:
            try:
                result = subprocess.run([exiftool, "-json", local_path], capture_output=True, text=True).stdout
                return json.loads(result)[0]
            except Exception:
                return None


class WavConverter(MediaConverter):
    """
    Converts WAV files to markdown via extraction of metadata (if `exiftool` is installed), and speech transcription (if `speech_recognition` is installed).
    """

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not a XLSX
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".wav":
            return None

        md_content = ""

        # Add metadata
        metadata = self._get_metadata(local_path)
        if metadata:
            for f in [
                "Title",
                "Artist",
                "Author",
                "Band",
                "Album",
                "Genre",
                "Track",
                "DateTimeOriginal",
                "CreateDate",
                "Duration",
            ]:
                if f in metadata:
                    md_content += f"{f}: {metadata[f]}\n"

        # Transcribe
        try:
            transcript = self._transcribe_audio(local_path)
            md_content += "\n\n### Audio Transcript:\n" + ("[No speech detected]" if transcript == "" else transcript)
        except Exception:
            md_content += "\n\n### Audio Transcript:\nError. Could not transcribe this audio."

        return DocumentConverterResult(
            title=None,
            text_content=md_content.strip(),
        )

    def _transcribe_audio(self, local_path) -> str:
        recognizer = sr.Recognizer()
        with sr.AudioFile(local_path) as source:
            audio = recognizer.record(source)
            return recognizer.recognize_google(audio).strip()


class Mp3Converter(WavConverter):
    """
    Converts MP3 and M4A files to markdown via extraction of metadata (if `exiftool` is installed), and speech transcription (if `speech_recognition` AND `pydub` are installed).
    """

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not a MP3
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".mp3", ".m4a"]:
            return None

        md_content = ""

        # Add metadata
        metadata = self._get_metadata(local_path)
        if metadata:
            for f in [
                "Title",
                "Artist",
                "Author",
                "Band",
                "Album",
                "Genre",
                "Track",
                "DateTimeOriginal",
                "CreateDate",
                "Duration",
            ]:
                if f in metadata:
                    md_content += f"{f}: {metadata[f]}\n"

        # Transcribe
        handle, temp_path = tempfile.mkstemp(suffix=".wav")
        os.close(handle)
        try:
            if extension.lower() == ".mp3":
                sound = pydub.AudioSegment.from_mp3(local_path)
            else:
                sound = pydub.AudioSegment.from_file(local_path, format="m4a")
            sound.export(temp_path, format="wav")

            _args = dict()
            _args.update(kwargs)
            _args["file_extension"] = ".wav"

            try:
                transcript = super()._transcribe_audio(temp_path).strip()
                md_content += "\n\n### Audio Transcript:\n" + (
                    "[No speech detected]" if transcript == "" else transcript
                )
            except Exception:
                md_content += "\n\n### Audio Transcript:\nError. Could not transcribe this audio."

        finally:
            os.unlink(temp_path)

        # Return the result
        return DocumentConverterResult(
            title=None,
            text_content=md_content.strip(),
        )


class ZipConverter(DocumentConverter):
    """
    Extracts ZIP files to a permanent local directory and returns a listing of extracted files.
    """

    def __init__(self, extract_dir: str = "downloads"):
        """
        Initialize with path to extraction directory.

        Args:
            extract_dir: The directory where files will be extracted. Defaults to "downloads"
        """
        self.extract_dir = extract_dir
        # Create the extraction directory if it doesn't exist
        os.makedirs(self.extract_dir, exist_ok=True)

    def convert(self, local_path: str, **kwargs: Any) -> Union[None, DocumentConverterResult]:
        # Bail if not a ZIP file
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".zip":
            return None

        # Verify it's actually a ZIP file
        if not zipfile.is_zipfile(local_path):
            return None

        # Extract all files and build list
        extracted_files = []
        with zipfile.ZipFile(local_path, "r") as zip_ref:
            # Extract all files
            zip_ref.extractall(self.extract_dir)
            # Get list of all files
            for file_path in zip_ref.namelist():
                # Skip directories
                if not file_path.endswith("/"):
                    extracted_files.append(self.extract_dir + "/" + file_path)

        # Sort files for consistent output
        extracted_files.sort()

        # Build the markdown content
        md_content = "Downloaded the following files:\n"
        for file in extracted_files:
            md_content += f"* {file}\n"

        return DocumentConverterResult(title="Extracted Files", text_content=md_content.strip())


class ImageConverter(MediaConverter):
    """
    Converts images to markdown via extraction of metadata (if `exiftool` is installed), OCR (if `easyocr` is installed), and description via a multimodal LLM (if an mlm_client is configured).
    """

    def convert(self, local_path, **kwargs: Any) -> Union[None, DocumentConverterResult]:
        # Bail if not a XLSX
        file_name = os.path.basename(local_path)
        # Split the file name into name and extension
        name, extension = os.path.splitext(file_name)
        if extension.lower() not in [".jpg", ".jpeg", ".png", "webp"]:
            # print("image format not supported!")
            return None

        md_content = ""

        # Add metadata
        metadata = self._get_metadata(local_path)
        if metadata:
            for f in [
                "ImageSize",
                "Title",
                "Caption",
                "Description",
                "Keywords",
                "Artist",
                "Author",
                "DateTimeOriginal",
                "CreateDate",
                "GPSPosition",
            ]:
                if f in metadata:
                    md_content += f"{f}: {metadata[f]}\n"

        # Try describing the image with GPTV
        md_content += (
            "\n# Description:\n"
            + self._get_mlm_description(local_path, extension, prompt=kwargs.get("mlm_prompt").strip())
            + "\n"
        )

        return DocumentConverterResult(
            title=name,
            text_content=md_content,
        )

    def _get_mlm_description(self, local_path, extension, prompt=None):
        if prompt is None or prompt.strip() == "":
            prompt = "Write a detailed description for this image."

        sys.stderr.write(f"MLM Prompt:\n{prompt}\n")

        data_uri = ""
        with open(local_path, "rb") as image_file:
            content_type, encoding = mimetypes.guess_type("_dummy" + extension)
            if content_type is None:
                content_type = "image/jpeg"
            image_base64 = base64.b64encode(image_file.read()).decode("utf-8")
            data_uri = f"data:{content_type};base64,{image_base64}"
        print('image send to vlm...')
        result = run_base64_analyser(prompt,data_uri,is_revised=False)
        return result


class FileConversionException(Exception):
    pass


class UnsupportedFormatException(Exception):
    pass


class MarkdownConverter:
    """(In preview) An extremely simple text-based document reader, suitable for LLM use.
    This reader will convert common file-types or webpages to Markdown."""

    def __init__(
        self,
        requests_session: Optional[requests.Session] = None,
        mlm_client: Optional[Any] = None,
        mlm_model: Optional[Any] = None,
    ):
        if requests_session is None:
            self._requests_session = requests.Session()
        else:
            self._requests_session = requests_session

        self._mlm_client = mlm_client
        self._mlm_model = mlm_model

        self._page_converters: List[DocumentConverter] = []

        # Register converters for successful browsing operations
        # Later registrations are tried first / take higher priority than earlier registrations
        # To this end, the most specific converters should appear below the most generic converters
        self.register_page_converter(PlainTextConverter())
        self.register_page_converter(HtmlConverter())
        self.register_page_converter(WikipediaConverter())
        self.register_page_converter(YouTubeConverter())
        self.register_page_converter(DocxConverter())
        self.register_page_converter(XlsxConverter())
        self.register_page_converter(PptxConverter())
        self.register_page_converter(WavConverter())
        self.register_page_converter(Mp3Converter())
        self.register_page_converter(ImageConverter())
        self.register_page_converter(ZipConverter())
        self.register_page_converter(PdfConverter())

    def convert(
        self, source: Union[str, requests.Response], **kwargs: Any
    ) -> DocumentConverterResult:  # TODO: deal with kwargs
        """
        Args:
            - source: can be a string representing a path or url, or a requests.response object
            - extension: specifies the file extension to use when interpreting the file. If None, infer from source (path, uri, content-type, etc.)
        """

        # Local path or url
        if isinstance(source, str):
            if source.startswith("http://") or source.startswith("https://") or source.startswith("file://"):
                return self.convert_url(source, **kwargs)
            else:
                return self.convert_local(source, **kwargs)
        # Request response
        elif isinstance(source, requests.Response):
            return self.convert_response(source, **kwargs)

    def convert_local(self, path: str, **kwargs: Any) -> DocumentConverterResult:  # TODO: deal with kwargs
        # Prepare a list of extensions to try (in order of priority)
        ext = kwargs.get("file_extension")
        extensions = [ext] if ext is not None else []

        # Get extension alternatives from the path and puremagic
        base, ext = os.path.splitext(path)
        # print(base)
        print('file ext is :', ext)
        self._append_ext(extensions, ext)
        self._append_ext(extensions, self._guess_ext_magic(path))

        # Convert
        return self._convert(path, extensions, **kwargs)

    # TODO what should stream's type be?
    def convert_stream(self, stream: Any, **kwargs: Any) -> DocumentConverterResult:  # TODO: deal with kwargs
        # Prepare a list of extensions to try (in order of priority)
        ext = kwargs.get("file_extension")
        extensions = [ext] if ext is not None else []

        # Save the file locally to a temporary file. It will be deleted before this method exits
        handle, temp_path = tempfile.mkstemp()
        fh = os.fdopen(handle, "wb")
        result = None
        try:
            # Write to the temporary file
            content = stream.read()
            if isinstance(content, str):
                fh.write(content.encode("utf-8"))
            else:
                fh.write(content)
            fh.close()

            # Use puremagic to check for more extension options
            self._append_ext(extensions, self._guess_ext_magic(temp_path))

            # Convert
            result = self._convert(temp_path, extensions, **kwargs)
        # Clean up
        finally:
            try:
                fh.close()
            except Exception:
                pass
            os.unlink(temp_path)

        return result

    def convert_url(self, url: str, **kwargs: Any) -> DocumentConverterResult:  # TODO: fix kwargs type
        # Send a HTTP request to the URL
        user_agent = "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/119.0.0.0 Safari/537.36 Edg/119.0.0.0"
        response = self._requests_session.get(url, stream=True, headers={"User-Agent": user_agent})
        response.raise_for_status()
        return self.convert_response(response, **kwargs)

    def convert_response(
        self, response: requests.Response, **kwargs: Any
    ) -> DocumentConverterResult:  # TODO fix kwargs type
        # Prepare a list of extensions to try (in order of priority)
        ext = kwargs.get("file_extension")
        print('ext: ', ext)
        extensions = [ext] if ext is not None else []

        # Guess from the mimetype
        content_type = response.headers.get("content-type", "").split(";")[0]
        self._append_ext(extensions, mimetypes.guess_extension(content_type))

        # Read the content disposition if there is one
        content_disposition = response.headers.get("content-disposition", "")
        m = re.search(r"filename=([^;]+)", content_disposition)
        if m:
            base, ext = os.path.splitext(m.group(1).strip("\"'"))
            self._append_ext(extensions, ext)

        # Read from the extension from the path
        base, ext = os.path.splitext(urlparse(response.url).path)
        self._append_ext(extensions, ext)

        # Save the file locally to a temporary file. It will be deleted before this method exits
        handle, temp_path = tempfile.mkstemp()
        fh = os.fdopen(handle, "wb")
        result = None
        try:
            # Download the file
            for chunk in response.iter_content(chunk_size=512):
                fh.write(chunk)
            fh.close()

            # Use puremagic to check for more extension options
            self._append_ext(extensions, self._guess_ext_magic(temp_path))

            # Convert
            result = self._convert(temp_path, extensions, url=response.url)
        except Exception as e:
            print(f"Error in converting: {e}")

        # Clean up
        finally:
            try:
                fh.close()
            except Exception:
                pass
            os.unlink(temp_path)

        return result

    def _convert(self, local_path: str, extensions: List[Union[str, None]], **kwargs) -> DocumentConverterResult:
        error_trace = ""
        for ext in extensions + [None]:  # Try last with no extension
            # print(ext)
            for converter in self._page_converters:
                _kwargs = copy.deepcopy(kwargs)

                # Overwrite file_extension appropriately
                if ext is None:
                    if "file_extension" in _kwargs:
                        del _kwargs["file_extension"]
                else:
                    _kwargs.update({"file_extension": ext})
                    # print(_kwargs)

                # Copy any additional global options
                if "mlm_client" not in _kwargs and self._mlm_client is not None:
                    _kwargs["mlm_client"] = self._mlm_client

                if "mlm_model" not in _kwargs and self._mlm_model is not None:
                    _kwargs["mlm_model"] = self._mlm_model

                # If we hit an error log it and keep trying
                try:
                    res = converter.convert(local_path, **_kwargs)
                except Exception:
                    error_trace = ("\n\n" + traceback.format_exc()).strip()

                if res is not None:
                    # Normalize the content
                    res.text_content = "\n".join([line.rstrip() for line in re.split(r"\r?\n", res.text_content)])
                    res.text_content = re.sub(r"\n{3,}", "\n\n", res.text_content)

                    # Todo
                    return res
        # If we got this far without success, report any exceptions
        if len(error_trace) > 0:
            raise FileConversionException(
                f"Could not convert '{local_path}' to Markdown. File type was recognized as {extensions}. While converting the file, the following error was encountered:\n\n{error_trace}"
            )

        # Nothing can handle it!
        raise UnsupportedFormatException(
            f"Could not convert '{local_path}' to Markdown. The formats {extensions} are not supported."
        )

    def _append_ext(self, extensions, ext):
        """Append a unique non-None, non-empty extension to a list of extensions."""
        if ext is None:
            return
        ext = ext.strip()
        if ext == "":
            return
        # if ext not in extensions:
        if True:
            extensions.append(ext)

    def _guess_ext_magic(self, path):
        """Use puremagic (a Python implementation of libmagic) to guess a file's extension based on the first few bytes."""
        # Use puremagic to guess
        try:
            guesses = puremagic.magic_file(path)
            if len(guesses) > 0:
                ext = guesses[0].extension.strip()
                if len(ext) > 0:
                    return ext
        except FileNotFoundError:
            pass
        except IsADirectoryError:
            pass
        except PermissionError:
            pass
        return None

    def register_page_converter(self, converter: DocumentConverter) -> None:
        """Register a page text converter."""
        self._page_converters.insert(0, converter)

# from marker.convert import convert_single_pdf
# import marker

if __name__ == "__main__":
    md_converter = MarkdownConverter()
    # result = md_converter.convert((r"G:\work\smolagent\gaia\063800f6-8832-4856-972b-17b877612533.png"), mlm_prompt="what's the number at the top left corner?")
    # result = md_converter.convert(r"G:\360Downloads\7245af7c-404e-4d60-9ef4-94ed301e5315.jpg",mlm_prompt="describe this image")
    result = md_converter.convert(r"G:\360Downloads\198ffd8f-6041-458d-bacc-fe49872cfa43.txt")
    # pdf_path_digital = r"G:\360Downloads\021a5339-744f-42b7-bd9b-9368b3efda7a.pdf" # "G:\360Downloads\Bei FE 20250517.pdf"
    # pdf_path_digital = r"G:\360Downloads\Bei FE 20250517.pdf"
    # result_digital = get_pdf_type(pdf_path_digital)
    # print(f"'{pdf_path_digital}': {result_digital['overall_type']} (Pages: {result_digital['page_types']})")
    print(result.text_content)
